# Construye la presentacion de la Clase 11 (estilo Tec) con python-pptx.
import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
from tec_template import (prs, BLANK, add_bg, add_header, add_text, add_card,
                          add_formula_box, add_footer, slide_portada,
                          slide_cierre, NAVY, LIGHT_BLUE, SOFT_BLUE,
                          SOFT_GREEN, DARK_GREEN, SOFT_RED, DARK_RED,
                          HIGHLIGHT_GREEN, GRAY, WHITE, BLACK)
from pptx.util import Inches, Pt

IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "img")
FOOTER = "CD II - Clase 11"
n = [0]

def nueva(titulo, subtitulo=None):
    n[0] += 1
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    add_header(slide, titulo)
    if subtitulo:
        add_text(slide, subtitulo, Inches(0.4), Inches(1.12), Inches(9.2),
                 Inches(0.45), size=14, italic=True)
    add_footer(slide, n[0], titulo_corto=FOOTER)
    return slide

def fuente(slide, texto, y=7.0):
    add_text(slide, texto, Inches(0.3), Inches(y), Inches(6.4), Inches(0.3),
             size=9, color=GRAY)

# ---------- 1. Portada ----------
n[0] += 1
slide_portada("Web scraping",
              "Clase 11 - Cuando el dato vive en una página y nadie lo publicó en CSV",
              "Ciencia de Datos II - 2 de septiembre de 2026")

# ---------- 2. Agenda ----------
slide = nueva("Hoy aprendemos a leer la web con código")
bloques = [
    ("1. HTML básico", "Etiquetas, atributos y selectores\nCSS: el mapa de toda página."),
    ("2. rvest", "read_html(), html_elements(),\nhtml_text2() y html_table()."),
    ("3. Ética y legalidad", "robots.txt, términos de uso y\ntasas de solicitud con Sys.sleep()."),
    ("4. La era de la IA", "Agentes que navegan, bloqueos\nanti-bot y licencias de contenido."),
]
for (x, y), (t, b) in zip([(0.3, 1.7), (5.05, 1.7), (0.3, 4.15), (5.05, 4.15)], bloques):
    add_card(slide, Inches(x), Inches(y), Inches(4.65), Inches(2.1),
             title=t, body=b, title_size=16, body_size=13, line_spacing=1.35)
add_card(slide, Inches(0.3), Inches(6.45), Inches(9.4), Inches(0.6),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Además: el equipo expositor presenta 'scraping en la era de la IA'.",
         title_size=12.5)

# ---------- 3. Repaso ----------
slide = nueva("Venimos de las fuentes ordenadas: datos abiertos")
rep = [
    ("El patrón data-raw/", "La descarga es parte del análisis: script que "
     "baja el archivo, README con URL y fecha, limpieza aparte."),
    ("download.file() y read_csv()", "Cuando hay URL directa a un CSV, dos "
     "funciones resuelven la ingesta completa."),
    ("¿Y si no hay CSV?", "Muchos datos valiosos solo existen como páginas "
     "web. Hoy aprendemos a extraerlos con respeto y con código."),
]
y = 1.7
for t, b in rep:
    add_card(slide, Inches(0.3), Inches(y), Inches(9.4), Inches(1.45),
             title=t, body=b, title_size=14, body_size=12.5)
    y += 1.62

# ---------- 4. Cuando scrapear ----------
slide = nueva("El scraping es el plan C: antes van datos y APIs")
jer = [
    ("Plan A: datos abiertos", "Si existe el CSV oficial, se descarga y "
     "listo (clase 10). Máxima confiabilidad, mínimo esfuerzo.",
     SOFT_GREEN, DARK_GREEN),
    ("Plan B: API oficial", "Si el sitio ofrece una API, es la vía "
     "diseñada para máquinas: estable y documentada (clase 12).",
     LIGHT_BLUE, NAVY),
    ("Plan C: scraping", "Solo cuando el dato existe únicamente como "
     "página web. Más frágil: si el sitio cambia su HTML, el script se "
     "rompe.", SOFT_RED, DARK_RED),
]
y = 1.7
for t, b, f, tc in jer:
    add_card(slide, Inches(0.3), Inches(y), Inches(9.4), Inches(1.5),
             fill=f, title_color=tc, title=t, body=b, title_size=14,
             body_size=12.5)
    y += 1.67

# ---------- 5. Arbol HTML ----------
slide = nueva("Toda página web es un árbol de etiquetas HTML")
slide.shapes.add_picture(os.path.join(IMG, "arbol_html.png"),
                         Inches(0.9), Inches(1.55), width=Inches(8.2))
add_card(slide, Inches(0.3), Inches(6.35), Inches(9.4), Inches(0.7),
         fill=LIGHT_BLUE,
         body="Para verlo en vivo: clic derecho en cualquier página, "
              "'Inspeccionar'. Ese árbol es lo que scrapearemos.",
         body_size=12)

# ---------- 6. Etiquetas y atributos ----------
slide = nueva("Etiquetas, atributos y clases: los ganchos")
add_card(slide, Inches(0.3), Inches(1.6), Inches(9.4), Inches(1.75),
         fill=SOFT_BLUE, title="Anatomía de un elemento",
         body='<a class="nota" href="https://ejemplo.mx">Leer más</a>\n\n'
              'a = etiqueta | class y href = atributos | Leer más = texto',
         title_size=14, body_size=13, line_spacing=1.35)
add_card(slide, Inches(0.3), Inches(3.6), Inches(4.65), Inches(2.9),
         title="Etiquetas frecuentes",
         body="h1...h6: títulos\np: párrafos\na: enlaces (href)\n"
              "table, tr, td: tablas\ndiv, span: contenedores\n"
              "img: imágenes (src)",
         title_size=14, body_size=12.5, line_spacing=1.35)
add_card(slide, Inches(5.05), Inches(3.6), Inches(4.65), Inches(2.9),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Clases e identificadores",
         body="class agrupa elementos del mismo tipo visual (muchos "
              "por página).\n\nid identifica un elemento único.\n\n"
              "Son los ganchos favoritos del scraping.",
         title_size=14, body_size=12.5, line_spacing=1.35)

# ---------- 7. Selectores CSS ----------
slide = nueva("Los selectores CSS apuntan al dato que queremos")
sel = [
    ('Por etiqueta:  "table"', "Todas las tablas de la página. Así se "
     "extraen tablas de Wikipedia.", LIGHT_BLUE, NAVY),
    ('Por clase:  ".archive-item-date"', "Todos los elementos con esa "
     "clase; el punto inicial es obligatorio. El caso típico: fechas o "
     "títulos de una lista de artículos.", LIGHT_BLUE, NAVY),
    ('Por id:  "#contenido"', "El elemento único con ese identificador; "
     "se antecede con el símbolo de gato.", LIGHT_BLUE, NAVY),
    ('Combinados:  "div.nota a"', "Los enlaces que viven dentro de un div "
     "con clase nota: se leen de izquierda a derecha, de padre a hijo.",
     SOFT_GREEN, DARK_GREEN),
]
y = 1.6
for t, b, f, tc in sel:
    add_card(slide, Inches(0.3), Inches(y), Inches(9.4), Inches(1.2),
             fill=f, title_color=tc, title=t, body=b, title_size=13,
             body_size=11.5, line_spacing=1.2)
    y += 1.33
fuente(slide, "El selector .archive-item-date proviene del ejercicio del curso anterior (atiempo.tv).")

# ---------- 8. rvest ----------
slide = nueva("rvest: el paquete de scraping del tidyverse")
add_card(slide, Inches(0.3), Inches(1.6), Inches(4.65), Inches(2.6),
         title="Qué es",
         body="Paquete de R inspirado en beautifulsoup (Python), diseñado "
              "para encadenarse con %>%.\n\nSe instala una vez: "
              "install.packages(\"rvest\").",
         title_size=14, body_size=12.5, line_spacing=1.35)
add_card(slide, Inches(5.05), Inches(1.6), Inches(4.65), Inches(2.6),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Sus cuatro verbos",
         body="read_html(): descarga el árbol\nhtml_elements(): selecciona "
              "ramas\nhtml_text2(): extrae el texto\nhtml_table(): extrae "
              "tablas",
         title_size=14, body_size=12.5, line_spacing=1.35)
add_card(slide, Inches(0.3), Inches(4.5), Inches(9.4), Inches(1.9),
         fill=HIGHLIGHT_GREEN, title="La lógica de siempre",
         body="Igual que con dplyr: pocas funciones que se combinan. "
              "Aprender rvest es aprender a escribir buenos selectores; el "
              "resto es el tidyverse que ya dominan.",
         title_size=14, body_size=13, line_spacing=1.4)

# ---------- 9. Flujo rvest ----------
slide = nueva("El flujo rvest: leer una vez, seleccionar, extraer")
slide.shapes.add_picture(os.path.join(IMG, "flujo_rvest.png"),
                         Inches(0.35), Inches(1.8), width=Inches(9.3))
add_card(slide, Inches(0.3), Inches(5.55), Inches(9.4), Inches(1.3),
         fill=LIGHT_BLUE,
         body="read_html() se llama UNA vez y el resultado se guarda en un "
              "objeto. Probar selectores sobre ese objeto no vuelve a "
              "molestar al sitio.", body_size=12.5, line_spacing=1.35)

# ---------- 10. read_html y html_elements ----------
slide = nueva("read_html() trae la página; html_elements() filtra")
add_card(slide, Inches(0.3), Inches(1.6), Inches(9.4), Inches(2.4),
         fill=SOFT_BLUE, title="Descargar y seleccionar",
         body='library(rvest)\n\n'
              'pagina <- read_html("https://es.wikipedia.org/wiki/Nayarit")\n\n'
              'titulos <- pagina %>%\n'
              '  html_elements("h2")',
         title_size=14, body_size=12.5, line_spacing=1.3),
add_card(slide, Inches(0.3), Inches(4.3), Inches(4.65), Inches(2.2),
         title="html_elements() acepta",
         body="Selectores CSS (etiqueta, .clase, #id) y también XPath "
              "para casos rebuscados. En singular, html_element() trae "
              "solo el primero.",
         title_size=13.5, body_size=12, line_spacing=1.3)
add_card(slide, Inches(5.05), Inches(4.3), Inches(4.65), Inches(2.2),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="El resultado aún no es tabla",
         body="Devuelve nodos del árbol. Para volverlos datos usables "
              "faltan los extractores de la siguiente diapositiva.",
         title_size=13.5, body_size=12, line_spacing=1.3)

# ---------- 11. html_text2 y html_attr ----------
slide = nueva("html_text2() y html_attr() extraen texto y enlaces")
add_card(slide, Inches(0.3), Inches(1.6), Inches(9.4), Inches(2.4),
         fill=SOFT_BLUE, title="Extraer contenido y atributos",
         body='titulos %>% html_text2()\n'
              '# "Historia" "Geografía" "Economía" ...\n\n'
              'pagina %>%\n'
              '  html_elements("a") %>%\n'
              '  html_attr("href")   # los enlaces de la página',
         title_size=14, body_size=12.5, line_spacing=1.3)
add_card(slide, Inches(0.3), Inches(4.3), Inches(4.65), Inches(2.2),
         title="¿Por qué text2 y no text?",
         body="html_text2() limpia espacios y saltos de línea como los "
              "ve el navegador; html_text() entrega el texto crudo. Casi "
              "siempre conviene text2.",
         title_size=13.5, body_size=12, line_spacing=1.3)
add_card(slide, Inches(5.05), Inches(4.3), Inches(4.65), Inches(2.2),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="html_attr() para metadatos",
         body="Los enlaces (href) y las imágenes (src) viven en "
              "atributos, no en el texto. Así se scrapean listas de "
              "artículos con su URL.",
         title_size=13.5, body_size=12, line_spacing=1.3)

# ---------- 12. html_table ----------
slide = nueva("html_table() convierte tablas HTML en tibbles")
add_card(slide, Inches(0.3), Inches(1.6), Inches(9.4), Inches(2.6),
         fill=SOFT_BLUE, title="El camino corto para tablas",
         body='tablas <- read_html(url_wiki) %>%\n'
              '  html_elements("table") %>%\n'
              '  html_table()   # lista de tibbles\n\n'
              'poblacion <- tablas[[2]] %>%    # la que nos interesa\n'
              '  janitor::clean_names()',
         title_size=14, body_size=12.5, line_spacing=1.3)
add_card(slide, Inches(0.3), Inches(4.5), Inches(9.4), Inches(1.9),
         fill=HIGHLIGHT_GREEN, title="Idea clave",
         body="Si el dato ya está en una tabla HTML, html_table() hace "
              "todo el trabajo estructural: solo queda elegir la tabla "
              "correcta de la lista y limpiar nombres y tipos.",
         title_size=14, body_size=13, line_spacing=1.4)

# ---------- 13. robots.txt ----------
slide = nueva("robots.txt es una convención, no un candado")
add_card(slide, Inches(0.3), Inches(1.6), Inches(4.65), Inches(2.9),
         title="Qué es",
         body="Archivo público en la raíz del sitio (ejemplo.mx/"
              "robots.txt) donde el dueño declara qué rutas pueden "
              "visitar los bots y cuáles no.\n\nSe revisa ANTES de "
              "scrapear; leerlo toma un minuto.",
         title_size=14, body_size=12, line_spacing=1.3)
add_card(slide, Inches(5.05), Inches(1.6), Inches(4.65), Inches(2.9),
         fill=SOFT_RED, title_color=DARK_RED,
         title="Qué NO es",
         body="No es un mecanismo técnico ni una ley: es una convención "
              "voluntaria.\n\nEn 2025 Cloudflare documentó que Perplexity "
              "accedía con rastreadores encubiertos a sitios que lo "
              "prohibían; perdió su estatus de bot verificado.",
         title_size=14, body_size=12, line_spacing=1.3)
add_card(slide, Inches(0.3), Inches(4.75), Inches(9.4), Inches(1.7),
         fill=HIGHLIGHT_GREEN, title="Nuestra regla en el curso",
         body="Si robots.txt o los términos de uso dicen que no, es no. "
              "La reputación (y la calificación) no se juegan por una "
              "tabla.", title_size=14, body_size=13, line_spacing=1.4)
fuente(slide, "Fuente: Cloudflare Blog, 4 de agosto de 2025 (consultado el 1 de agosto de 2026).")

# ---------- 14. Semaforo ----------
slide = nueva("El semáforo ético del scraping")
slide.shapes.add_picture(os.path.join(IMG, "semaforo_etico.png"),
                         Inches(0.8), Inches(1.7), width=Inches(8.4))
add_card(slide, Inches(0.3), Inches(6.15), Inches(9.4), Inches(0.75),
         fill=LIGHT_BLUE,
         body="Ante la duda: buscar la fuente alternativa o pedir permiso "
              "al sitio. Casi siempre contestan.", body_size=12)

# ---------- 15. Cortesia tecnica ----------
slide = nueva("Scrapear despacio y con identificación honesta")
add_card(slide, Inches(0.3), Inches(1.6), Inches(9.4), Inches(2.3),
         fill=SOFT_BLUE, title="Pausas entre solicitudes",
         body='urls_paginas %>%\n'
              '  lapply(function(u) {\n'
              '    Sys.sleep(2)          # dos segundos de cortesía\n'
              '    read_html(u)\n'
              '  })',
         title_size=14, body_size=12.5, line_spacing=1.3)
add_card(slide, Inches(0.3), Inches(4.2), Inches(4.65), Inches(2.3),
         title="Por qué importa",
         body="Cientos de solicitudes por segundo pueden tirar un sitio "
              "chico: eso ya no es recolección, es daño. Las pausas y "
              "los horarios valle son la diferencia.",
         title_size=13.5, body_size=12, line_spacing=1.3)
add_card(slide, Inches(5.05), Inches(4.2), Inches(4.65), Inches(2.3),
         fill=SOFT_GREEN, title_color=DARK_GREEN,
         title="Identificarse",
         body="Un agente de usuario honesto con correo de contacto, no "
              "disfrazarse de navegador. Lección del caso hiQ-LinkedIn: "
              "lo público no es delito, pero el contrato del sitio sí "
              "obliga.",
         title_size=13.5, body_size=12, line_spacing=1.3)
fuente(slide, "Caso hiQ Labs v. LinkedIn (2019-2022): el acceso público no viola la CFAA; los términos de uso sí obligan.")

# ---------- 16. Bots dominan el trafico ----------
slide = nueva("Los bots ya generan más tráfico que los humanos")
slide.shapes.add_picture(os.path.join(IMG, "rastreo_referencias.png"),
                         Inches(1.0), Inches(1.6), width=Inches(8.0))
add_card(slide, Inches(0.3), Inches(5.85), Inches(9.4), Inches(1.05),
         fill=LIGHT_BLUE,
         body="En junio de 2026 los bots generaron 57.5% del tráfico HTML "
              "(los humanos, 42.5%). El trato implícito de la web (te "
              "rastreo y te mando visitas) se rompió: los bots de IA leen "
              "miles de páginas y regresan casi nada.",
         body_size=12, line_spacing=1.3)
fuente(slide, "Fuente: compilación de Digital Applied sobre Cloudflare Radar (consultada el 1 de agosto de 2026).", y=7.05)

# ---------- 17. La web responde ----------
slide = nueva("La web responde: bloqueo por defecto y peajes")
resp = [
    ("Bloqueo por defecto", "Desde julio de 2025 Cloudflare bloquea a los "
     "rastreadores de IA salvo permiso explícito: el modelo pasó de "
     "opt-out a opt-in, con el respaldo de más de 50 medios y "
     "plataformas.", SOFT_RED, DARK_RED),
    ("Pago por rastreo (pay per crawl)", "El sitio cobra por solicitud "
     "usando el código HTTP 402 (Payment Required); el bot declara cuánto "
     "está dispuesto a pagar. En beta desde 2025.", LIGHT_BLUE, NAVY),
    ("Qué significa para ustedes", "El scraping casual de sitios grandes "
     "será cada vez más difícil: más razón para preferir datos abiertos y "
     "APIs, y para scrapear solo donde son bienvenidos.", SOFT_GREEN,
     DARK_GREEN),
]
y = 1.7
for t, b, f, tc in resp:
    add_card(slide, Inches(0.3), Inches(y), Inches(9.4), Inches(1.5),
             fill=f, title_color=tc, title=t, body=b, title_size=14,
             body_size=12)
    y += 1.67
fuente(slide, "Fuente: comunicados y blog técnico de Cloudflare, julio de 2025 (consultados el 1 de agosto de 2026).")

# ---------- 18. Litigios y licencias ----------
slide = nueva("Entrenar con contenido ajeno ya se litiga y licencia")
lit = [
    ("NYT contra OpenAI y Microsoft", "Demanda activa desde 2023 por "
     "entrenar con millones de artículos; en 2026 sigue en etapa de "
     "pruebas, sin fallo de fondo.", SOFT_RED, DARK_RED),
    ("Acuerdo Anthropic - autores (2026)", "1,500 millones de dólares por "
     "haber usado libros de bibliotecas piratas: el mayor acuerdo de "
     "derechos de autor en EE. UU. Entrenar con obras compradas sí fue "
     "declarado uso legítimo.", SOFT_RED, DARK_RED),
    ("La vía del contrato", "AP, The Guardian, Washington Post y decenas "
     "más ya licencian su contenido a laboratorios de IA: el dato con "
     "valor se vende, no se regala.", LIGHT_BLUE, NAVY),
]
y = 1.7
for t, b, f, tc in lit:
    add_card(slide, Inches(0.3), Inches(y), Inches(9.4), Inches(1.5),
             fill=f, title_color=tc, title=t, body=b, title_size=14,
             body_size=12)
    y += 1.67
fuente(slide, "Fuentes: Authors Guild (20 de julio de 2026); Wikipedia (NYT v. OpenAI); Digiday 2025.")

# ---------- 19. Exposicion ----------
slide = nueva("El equipo expositor: scraping en la era de la IA",
              "10 minutos + preguntas.")
add_card(slide, Inches(0.3), Inches(1.9), Inches(9.4), Inches(2.2),
         title="Tema de hoy: web scraping en la nueva realidad de la IA",
         body="El equipo expositor trabaja con el documento base del curso "
              "(carpeta 13: agentes que navegan, bloqueos anti-bot, "
              "licenciamiento) y al menos una fuente propia. El resto del "
              "grupo: una pregunta por equipo.",
         title_size=15, body_size=13, line_spacing=1.4)
add_card(slide, Inches(0.3), Inches(4.3), Inches(9.4), Inches(1.6),
         fill=LIGHT_BLUE, title="Guía para escuchar",
         body="Mientras exponen, respondan en su cuaderno: si un agente "
              "de IA navega por ti, ¿quién es responsable de respetar "
              "robots.txt? ¿Cambia algo si el sitio cobra por rastreo?",
         title_size=14, body_size=12.5, line_spacing=1.35)

# ---------- 20. Ejercicio ----------
slide = nueva("Ejercicio: scrapear una tabla real de Wikipedia",
              "En parejas, 25 minutos; requiere laptop con R y rvest instalado.")
add_card(slide, Inches(0.3), Inches(1.8), Inches(9.4), Inches(3.0),
         title="Instrucciones",
         body="1. Elijan un artículo de Wikipedia en español con una "
              "tabla que les interese (población por estado, medallero, "
              "elecciones...).\n"
              "2. Revisen es.wikipedia.org/robots.txt: ¿scrapear "
              "artículos está permitido?\n"
              "3. Con read_html() + html_elements(\"table\") + "
              "html_table(), extraigan la lista de tablas y localicen la "
              "suya.\n"
              "4. Límpienla: clean_names(), tipos correctos con mutate() "
              "y un count() o summarise() que diga algo.\n"
              "5. Documenten en comentarios: URL, fecha de scrapeo y "
              "número de tabla.",
         title_size=15, body_size=12.5, line_spacing=1.45)
add_card(slide, Inches(0.3), Inches(5.0), Inches(9.4), Inches(1.4),
         fill=HIGHLIGHT_GREEN, title="Entregable (fin de la clase)",
         body="El script .R con el flujo completo y una captura del "
              "tibble final. Se sube al hilo de Canvas de la clase 11.",
         title_size=14, body_size=12.5)

# ---------- 21. Conclusiones ----------
slide = nueva("Conclusiones de la clase 11")
concl = (
    "1. Toda página es un árbol de etiquetas; los selectores CSS son la "
    "dirección de cada dato.\n"
    "2. rvest resuelve el flujo completo: read_html() para leer, "
    "html_elements() para seleccionar, html_text2() y html_table() para "
    "extraer.\n"
    "3. El scraping es el plan C: primero datos abiertos, luego APIs, al "
    "final HTML.\n"
    "4. robots.txt y los términos de uso se respetan; se scrapea "
    "despacio (Sys.sleep) y con identificación honesta.\n"
    "5. La era de la IA endureció la web: bloqueos por defecto, pago por "
    "rastreo y litigios millonarios. Scrapear bien es también saber "
    "cuándo no scrapear."
)
add_card(slide, Inches(0.3), Inches(1.7), Inches(9.4), Inches(4.9),
         fill=HIGHLIGHT_GREEN, title="Lo que se llevan hoy",
         body=concl, title_size=16, body_size=13, line_spacing=1.55)

# ---------- 22. Recursos ----------
slide = nueva("Recursos para profundizar")
rec = (
    "Documentación de rvest: rvest.tidyverse.org (incluye el tutorial "
    "'SelectorGadget' para descubrir selectores sin leer HTML).\n\n"
    "Wickham, H., Çetinkaya-Rundel, M. y Grolemund, G.: R for Data "
    "Science (2a ed.), capítulo de web scraping (r4ds.hadley.nz).\n\n"
    "Documento del curso (Canvas): 'Web scraping en la nueva realidad de "
    "la IA', con las 15 fuentes verificadas de esta clase.\n\n"
    "Para explorar HTML en vivo: el inspector del navegador (clic "
    "derecho, 'Inspeccionar') y validator.w3.org.\n\n"
    "Próxima clase: APIs con httr2, RAG y MCP; recuerden que el "
    "mini-proyecto 1 se entrega el 9 de septiembre."
)
add_card(slide, Inches(0.3), Inches(1.7), Inches(9.4), Inches(4.6),
         body=rec, body_size=13, line_spacing=1.5)

# ---------- 23. Cierre ----------
n[0] += 1
slide_cierre()

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Clase09_WebScraping.pptx")
prs.save(out)
print("Guardado:", out, "| slides:", len(prs.slides._sldIdLst))
